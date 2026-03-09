import io
from fastapi import APIRouter, UploadFile, File, Form, BackgroundTasks, Depends, HTTPException
from PyPDF2 import PdfReader
from pptx import Presentation

from app.config import get_settings, Settings
from app.embeddings.indexer import index_document
from app.utils.logger import get_logger

logger = get_logger("ingest")

router = APIRouter()

VALID_FILE_TYPES = [
    "medical_protocol",
    "lab_interpretation",
    "diagnostic_report",
    "ip_material",
    "other",
]


@router.post("/ingest")
async def ingest_document_endpoint(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    file_type: str = Form(...),
    user_id: str = Form(...),  # Passed from Next.js after JWT validation
    settings: Settings = Depends(get_settings),
):
    """
    Ingest a document into the vector store.

    Process:
    1. Extract text from PDF/document
    2. Chunk text semantically
    3. Generate embeddings
    4. Store in pgvector
    """
    # Validate file type
    if file_type not in VALID_FILE_TYPES:
        raise HTTPException(
            status_code=400,
            detail=f"Invalid file_type. Must be one of: {VALID_FILE_TYPES}",
        )

    # Read file content
    file_content = await file.read()
    mime_type = file.content_type or "application/octet-stream"

    # Extract text based on file type
    try:
        text_content = extract_text(file_content, mime_type)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to extract text: {str(e)}")

    if not text_content.strip():
        raise HTTPException(status_code=400, detail="No text content found in document")

    # Process in background for large documents
    background_tasks.add_task(
        process_document_task,
        user_id=user_id,
        filename=file.filename or "document",
        file_type=file_type,
        mime_type=mime_type,
        text_content=text_content,
    )

    return {
        "status": "processing",
        "filename": file.filename,
        "file_type": file_type,
        "message": "Document is being processed and indexed",
    }


@router.post("/ingest/sync")
async def ingest_document_sync(
    file: UploadFile = File(...),
    file_type: str = Form(...),
    user_id: str = Form(...),
    settings: Settings = Depends(get_settings),
):
    """
    Synchronous document ingestion - waits for processing to complete.

    Use this for smaller documents or when you need immediate confirmation.
    """
    if file_type not in VALID_FILE_TYPES:
        raise HTTPException(
            status_code=400,
            detail=f"Invalid file_type. Must be one of: {VALID_FILE_TYPES}",
        )

    file_content = await file.read()
    mime_type = file.content_type or "application/octet-stream"

    try:
        text_content = extract_text(file_content, mime_type)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to extract text: {str(e)}")

    if not text_content.strip():
        raise HTTPException(status_code=400, detail="No text content found in document")

    # Process synchronously
    result = await index_document(
        user_id=user_id,
        filename=file.filename or "document",
        file_type=file_type,
        mime_type=mime_type,
        text_content=text_content,
    )

    if result.get("status") == "error":
        raise HTTPException(status_code=500, detail=result.get("error", "Indexing failed"))

    return result


def extract_text(content: bytes, mime_type: str) -> str:
    """Extract text content from various file formats."""

    if mime_type == "application/pdf":
        return extract_pdf_text(content)

    elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return extract_pptx_text(content)

    elif mime_type in ["text/plain", "text/markdown"]:
        return content.decode("utf-8", errors="ignore")

    elif mime_type in ["application/json"]:
        return content.decode("utf-8", errors="ignore")

    else:
        # Try to decode as text
        try:
            return content.decode("utf-8", errors="ignore")
        except Exception:
            raise ValueError(f"Unsupported file type: {mime_type}")


def extract_pdf_text(content: bytes) -> str:
    """Extract text from PDF content."""
    reader = PdfReader(io.BytesIO(content))
    text_parts = []

    for page in reader.pages:
        text = page.extract_text()
        if text:
            text_parts.append(text)

    return "\n\n".join(text_parts)


def extract_pptx_text(content: bytes) -> str:
    """Extract text from PPTX content."""
    prs = Presentation(io.BytesIO(content))
    text_parts = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts = [f"## Slide {slide_num}"]

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)

            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip(" |"):
                        slide_texts.append(row_text)

        # Extract speaker notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_texts.append(f"\n**Notes:** {notes_text}")

        if len(slide_texts) > 1:
            text_parts.append("\n".join(slide_texts))

    return "\n\n---\n\n".join(text_parts)


async def process_document_task(
    user_id: str,
    filename: str,
    file_type: str,
    mime_type: str,
    text_content: str,
):
    """Background task to process and index a document."""
    try:
        await index_document(
            user_id=user_id,
            filename=filename,
            file_type=file_type,
            mime_type=mime_type,
            text_content=text_content,
        )
    except Exception as e:
        logger.error(f"Error indexing document {filename}: {e}")
