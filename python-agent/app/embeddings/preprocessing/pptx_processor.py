"""
PPTX Processor - Extract text from PowerPoint presentations.

Handles PPTX files by extracting:
- Slide text (shapes, text frames)
- Speaker notes
- Table content
- Grouped shapes

Vision-enhanced mode (extract_pptx_with_vision):
- Renders slides to PNG via LibreOffice + PyMuPDF
- Sends each slide image to Anthropic vision model
- Captures diagrams, charts, SmartArt that text-only extraction misses
"""

import base64
import subprocess
import tempfile
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches


def extract_pptx_text(pptx_path: Path) -> dict:
    """
    Extract text content from a PPTX file.

    Args:
        pptx_path: Path to the PPTX file

    Returns:
        Dictionary with:
        - text_content: Extracted text as markdown
        - slide_count: Number of slides
        - has_notes: Whether speaker notes were found
        - has_content: Whether meaningful content was found
    """
    prs = Presentation(str(pptx_path))

    content_parts = []
    slide_count = 0
    has_notes = False

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_count += 1
        slide_content = []

        # Add slide header
        slide_content.append(f"## Slide {slide_num}")

        # Extract text from all shapes
        for shape in slide.shapes:
            shape_text = _extract_shape_text(shape)
            if shape_text:
                slide_content.append(shape_text)

        # Extract speaker notes
        notes_text = _extract_notes(slide)
        if notes_text:
            has_notes = True
            slide_content.append(f"\n**Speaker Notes:**\n{notes_text}")

        # Only add slide if it has content beyond the header
        if len(slide_content) > 1:
            content_parts.append("\n\n".join(slide_content))

    full_text = "\n\n---\n\n".join(content_parts)

    return {
        "text_content": full_text,
        "slide_count": slide_count,
        "has_notes": has_notes,
        "has_content": len(full_text.strip()) > 50,
    }


SLIDE_VISION_PROMPT = (
    "Extract ALL text, protocol names, supplement names, frequencies, clinical data, "
    "chart values, and diagram information from this BFM seminar slide. Include any "
    "visible headings, bullet points, table data, and annotations. Format as structured "
    "markdown."
)


def _render_pptx_to_pngs(pptx_path: Path, output_dir: Path) -> list[Path]:
    """
    Render a PPTX file to one PNG per slide using LibreOffice + PyMuPDF.

    LibreOffice converts PPTX → PDF, then PyMuPDF rasterises each page.

    Args:
        pptx_path: Path to the PPTX file
        output_dir: Directory to write PNGs into

    Returns:
        Sorted list of PNG file paths (one per slide)
    """
    import fitz  # PyMuPDF
    import shutil

    # Step 1: PPTX → PDF via LibreOffice headless
    # Resolve the LibreOffice binary (PATH or macOS app bundle)
    lo_bin = shutil.which("libreoffice")
    if lo_bin is None:
        macos_path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
        if Path(macos_path).exists():
            lo_bin = macos_path
        else:
            raise FileNotFoundError(
                "LibreOffice not found. Install with: brew install --cask libreoffice"
            )

    subprocess.run(
        [
            lo_bin,
            "--headless",
            "--convert-to", "pdf",
            str(pptx_path),
            "--outdir", str(output_dir),
        ],
        check=True,
        capture_output=True,
        timeout=120,
    )

    pdf_path = output_dir / f"{pptx_path.stem}.pdf"
    if not pdf_path.exists():
        raise FileNotFoundError(f"LibreOffice did not produce expected PDF: {pdf_path}")

    # Step 2: PDF → PNGs via PyMuPDF (one image per page/slide)
    doc = fitz.open(str(pdf_path))
    png_paths: list[Path] = []
    for page_num in range(len(doc)):
        page = doc[page_num]
        # 2x zoom for high-res rendering (good balance of quality vs size)
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        png_path = output_dir / f"slide_{page_num + 1:03d}.png"
        pix.save(str(png_path))
        png_paths.append(png_path)
    doc.close()

    return sorted(png_paths)


def _vision_extract_slide(png_path: Path) -> str:
    """
    Send a single slide PNG to the vision model and return extracted markdown.

    Uses the Anthropic SDK with base64 image input.
    """
    from app.services.ai_client import get_sync_client, get_vision_model

    with open(png_path, "rb") as f:
        image_data = base64.b64encode(f.read()).decode("utf-8")

    client = get_sync_client()
    response = client.messages.create(
        model=get_vision_model(),
        max_tokens=2000,
        messages=[{
            "role": "user",
            "content": [
                {
                    "type": "image",
                    "source": {
                        "type": "base64",
                        "media_type": "image/png",
                        "data": image_data,
                    },
                },
                {"type": "text", "text": SLIDE_VISION_PROMPT},
            ],
        }],
    )
    return response.content[0].text


def extract_pptx_with_vision(pptx_path: Path) -> dict:
    """
    Extract content from a PPTX using LibreOffice rendering + vision API.

    Renders each slide to a PNG, sends to the vision model for comprehensive
    content extraction (captures diagrams, charts, SmartArt), then combines
    with native speaker notes.

    Args:
        pptx_path: Path to the PPTX file

    Returns:
        Dictionary with:
        - text_content: Vision-extracted text as markdown
        - slide_count: Number of slides
        - has_notes: Whether speaker notes were found
        - has_content: Whether meaningful content was found
    """
    prs = Presentation(str(pptx_path))
    slide_count = len(prs.slides)

    # Collect speaker notes from native python-pptx (not visible in rendered images)
    notes_by_slide: dict[int, str] = {}
    has_notes = False
    for slide_num, slide in enumerate(prs.slides, start=1):
        notes_text = _extract_notes(slide)
        if notes_text:
            notes_by_slide[slide_num] = notes_text
            has_notes = True

    # Render slides to PNGs and extract via vision
    content_parts: list[str] = []
    with tempfile.TemporaryDirectory() as tmpdir:
        tmp_path = Path(tmpdir)
        print(f"  Rendering {slide_count} slides via LibreOffice...")
        png_paths = _render_pptx_to_pngs(pptx_path, tmp_path)
        print(f"  Rendered {len(png_paths)} slide images")

        for i, png_path in enumerate(png_paths, start=1):
            slide_num = i
            print(f"  Vision extracting slide {slide_num}/{len(png_paths)}...")

            try:
                vision_text = _vision_extract_slide(png_path)
            except Exception as e:
                print(f"    Warning: Vision failed for slide {slide_num}: {e}")
                vision_text = ""

            # Build per-slide content
            parts = [f"## Slide {slide_num}"]
            if vision_text.strip():
                parts.append(vision_text.strip())

            notes = notes_by_slide.get(slide_num)
            if notes:
                parts.append(f"\n**Speaker Notes:**\n{notes}")

            # Only include slides that have some content
            if len(parts) > 1:
                content_parts.append("\n\n".join(parts))

    full_text = "\n\n---\n\n".join(content_parts)

    return {
        "text_content": full_text,
        "slide_count": slide_count,
        "has_notes": has_notes,
        "has_content": len(full_text.strip()) > 50,
    }


def _extract_shape_text(shape) -> Optional[str]:
    """Extract text from a shape, handling various shape types."""
    # Handle grouped shapes recursively
    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        group_texts = []
        for child_shape in shape.shapes:
            child_text = _extract_shape_text(child_shape)
            if child_text:
                group_texts.append(child_text)
        return "\n".join(group_texts) if group_texts else None

    # Handle tables
    if shape.has_table:
        return _table_to_markdown(shape.table)

    # Handle text frames
    if shape.has_text_frame:
        return _extract_text_frame(shape.text_frame)

    return None


def _extract_text_frame(text_frame) -> Optional[str]:
    """Extract text from a text frame, preserving paragraph structure."""
    paragraphs = []

    for para in text_frame.paragraphs:
        text = para.text.strip()
        if text:
            # Check indentation level for list items
            level = para.level if para.level else 0
            if level > 0:
                indent = "  " * level
                text = f"{indent}- {text}"
            paragraphs.append(text)

    return "\n".join(paragraphs) if paragraphs else None


def _extract_notes(slide) -> Optional[str]:
    """Extract speaker notes from a slide."""
    if not slide.has_notes_slide:
        return None

    notes_slide = slide.notes_slide
    notes_frame = notes_slide.notes_text_frame

    if notes_frame:
        notes_text = notes_frame.text.strip()
        return notes_text if notes_text else None

    return None


def _table_to_markdown(table) -> Optional[str]:
    """Convert a PPTX table to markdown format."""
    rows = []

    for row in table.rows:
        cells = []
        for cell in row.cells:
            cell_text = cell.text.strip().replace("\n", " ")
            cells.append(cell_text)
        if any(cells):  # Skip completely empty rows
            rows.append(cells)

    if not rows:
        return None

    # Build markdown table
    md_lines = []

    # Header row
    header = rows[0]
    md_lines.append("| " + " | ".join(header) + " |")
    md_lines.append("| " + " | ".join(["---"] * len(header)) + " |")

    # Data rows
    for row in rows[1:]:
        # Pad row to match header length
        while len(row) < len(header):
            row.append("")
        md_lines.append("| " + " | ".join(row) + " |")

    return "\n".join(md_lines)
